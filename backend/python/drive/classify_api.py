import json
import numpy as np
from flask import Flask, request, jsonify
from collections import defaultdict
from PyPDF2 import PdfReader
from docx import Document
from utils import embed
from pptx import Presentation
from PIL import Image
import easyocr
import cv2
import re
from textblob import TextBlob

import ssl
ssl._create_default_https_context = ssl._create_unverified_context

app = Flask(__name__)

reader = easyocr.Reader(['en'], gpu=False)

with open("syllabus_embeddings.json", "r", encoding="utf-8") as f:
    SYLLABUS = json.load(f)

SUBJECT_EMB = defaultdict(list)

for row in SYLLABUS:
    SUBJECT_EMB[row["subject"]].append(np.array(row["embedding"]))

SUBJECT_EMB = {
    s: np.mean(v, axis=0)
    for s, v in SUBJECT_EMB.items()
}

def cosine(a, b):
    return float(np.dot(a, b))

def chunk_text(text, max_words=300):
    words = text.split()
    return [{"text": " ".join(words[i:i+max_words])}
            for i in range(0, len(words), max_words)]

def clean_ocr_text(text):
    text = text.lower()

    # Fix common OCR confusions
    replacements = {
        "0": "o",
        "1": "l",
        "_": " ",
        "~": " ",
        "|": "l"
    }

    for k, v in replacements.items():
        text = text.replace(k, v)

    # Remove non-useful symbols
    text = re.sub(r"[^a-zA-Z0-9\s]", " ", text)

    # Remove extra spaces
    text = re.sub(r"\s+", " ", text).strip()

    return text

def correct_spelling(text):
    corrected_text = str(TextBlob(text).correct())
    return corrected_text

def extract_text_from_image(file):
    image = Image.open(file).convert("RGB")
    img = np.array(image)

    gray = cv2.cvtColor(img, cv2.COLOR_RGB2GRAY)
    gray = cv2.medianBlur(gray, 3)

    thresh = cv2.adaptiveThreshold(
        gray,
        255,
        cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
        cv2.THRESH_BINARY,
        31,
        2
    )


    results = reader.readtext(img)  

    texts = []

    for item in results:
        if len(item) == 3:
            bbox, text, _ = item
            texts.append(text)
        else:  # fallback in case EasyOCR returns unexpected structure
            texts.append(str(item))
    final_text = " ".join(texts)
    print("Extracted text from image:", final_text)
    return final_text


def extract_text(file):
    name = file.filename.lower()

    if name.endswith(".txt"):
        return file.read().decode("utf-8", errors="ignore")

    if name.endswith(".pdf"):
        pdf_reader = PdfReader(file)
        return "\n".join(p.extract_text() or "" for p in pdf_reader.pages)


    if name.endswith(".docx"):
        doc = Document(file)
        return "\n".join(p.text for p in doc.paragraphs)
    
    if name.endswith(".pptx"):
        prs = Presentation(file)
        text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)

        return "\n".join(text)

    if name.endswith((".jpg", ".jpeg", ".png")):
        raw_text = extract_text_from_image(file)
        cleaned_text = clean_ocr_text(raw_text)
        text = correct_spelling(cleaned_text)
        print("Cleaned OCR text:", text)
        return text

    return None

def get_top_k_chunks(chunks, subject, k=3):
    scored = []

    subject_emb = SUBJECT_EMB[subject]

    for c in chunks:
        emb = embed(c["text"])
        score = cosine(emb, subject_emb)
        scored.append((score, c))

    scored.sort(reverse=True, key=lambda x: x[0])
    return [c for _, c in scored[:k]]

def classify_subject(chunks):
    scores = defaultdict(list)

    for c in chunks:
        emb = embed(c["text"])
        for subject, s_emb in SUBJECT_EMB.items():
            scores[subject].append(cosine(emb, s_emb))

    avg = {k: np.mean(v) for k, v in scores.items()}
    best = max(avg, key=avg.get)

    # print("Subject scores:", avg)

    return best if avg[best] > 0.15 else None

def classify_unit(subject, chunks):
    top_chunks = get_top_k_chunks(chunks, subject, k=3)

    scores = defaultdict(list)

    for c in top_chunks:
        emb = embed(c["text"])

        for row in SYLLABUS:
            if row["subject"] != subject:
                continue

            key = (row["unit"], row["title"])
            scores[key].append(cosine(emb, row["embedding"]))

    if not scores:
        return None

    avg = {k: np.mean(v) for k, v in scores.items()}
    best = max(avg, key=avg.get)

    return {
        "unit": best[0],
        "title": best[1],
        "confidence": round(avg[best], 3)
    }


@app.route("/classify", methods=["POST"])
def classify():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    text = extract_text(request.files["file"])
    if not text:
        return jsonify({"error": "Empty file"}), 400

    chunks = chunk_text(text)

    subject = classify_subject(chunks)
    if subject is None:
        return jsonify({"subject": None, "unit": None})

    unit = classify_unit(subject, chunks)

    return jsonify({
        "subject": subject,
        "unit": unit["unit"] if unit else None,
        "unit_title": unit["title"] if unit else None,
        "confidence": unit["confidence"] if unit else None
    })

if __name__ == "__main__":
    app.run(port=8000, debug=True)

